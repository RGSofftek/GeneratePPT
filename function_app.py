import azure.functions as func
import logging
import os
import json
import time
from datetime import datetime
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from io import BytesIO
from azure.storage.fileshare import ShareServiceClient
import pandas as pd
import matplotlib.pyplot as plt
import tempfile
from PIL import Image
from azure.ai.inference import ChatCompletionsClient
from azure.core.credentials import AzureKeyCredential

# Configuración de variables de entorno
STORAGE_ACCOUNT_NAME = os.environ.get("STORAGE_ACCOUNT_NAME")
SAS_TOKEN = os.environ.get("SAS_TOKEN")
AUTH_TOKEN_EXPECTED = os.environ.get("AUTH_TOKEN")
FILE_SHARE_NAME = os.environ.get("FILE_SHARE_NAME")
DIRECTORY_NAME = os.environ.get("DIRECTORY_NAME")
TEMPLATES_DIRECTORY = f"{DIRECTORY_NAME}/templates"
INPUTS_DIRECTORY = f"{DIRECTORY_NAME}/inputs"
OUTPUTS_DIRECTORY = f"{DIRECTORY_NAME}/outputs"

# Validación de variables de entorno para Azure OpenAI
api_key = os.environ.get("AZURE_OPENAI_KEY")
if not api_key:
    raise Exception("A key should be provided to invoke the endpoint")

base_endpoint = os.environ.get("AZURE_OPENAI_ENDPOINT")
if not base_endpoint:
    raise Exception("An endpoint should be provided to invoke the service")

deployment_name = os.environ.get("AZURE_OPENAI_DEPLOYMENT", "gpt-4")
if not deployment_name:
    raise Exception("A deployment name should be provided")

# Asegurarse de que el endpoint base no tenga una barra doble o falte una barra
base_endpoint = base_endpoint.rstrip("/")
endpoint = f"{base_endpoint}/openai/deployments/{deployment_name}"

client = ChatCompletionsClient(
    endpoint=endpoint,
    credential=AzureKeyCredential(api_key)
)

# Nombres de archivo constantes
TMD_FILE = "TMD.xlsx"
USERS_FILE = "base_equipo.xlsx"
PASES_FILE = "Calidad_pases.xlsx"
REVISIONES_FILE = "Reversiones.xlsx"
MATURITY_LEVEL_FILE = "NIVEL_MADUREZ.xlsx"

service_url = f"https://{STORAGE_ACCOUNT_NAME}.file.core.windows.net"
SERVICE_CLIENT = ShareServiceClient(account_url=service_url, credential=SAS_TOKEN)
SHARE_CLIENT = SERVICE_CLIENT.get_share_client(FILE_SHARE_NAME)

def download_file_from_share(temp_dir: str, filename: str, folder: str) -> str:
    directory_client = SHARE_CLIENT.get_directory_client(folder)
    file_client = directory_client.get_file_client(filename)
    stream = file_client.download_file()
    local_file_path = os.path.join(temp_dir, filename)
    with open(local_file_path, 'wb') as f:
        f.write(stream.readall())
    logging.info(f"Descargado '{filename}' desde '{folder}' a '{local_file_path}'")
    return local_file_path

def upload_file_to_share(local_file_path: str, filename: str, folder: str) -> None:
    directory_client = SHARE_CLIENT.get_directory_client(folder)
    file_client = directory_client.get_file_client(filename)
    try:
        file_client.delete_file()
    except Exception:
        pass
    with open(local_file_path, 'rb') as data:
        file_client.upload_file(data)
    logging.info(f"Subido archivo '{filename}' a '{folder}'")

def generate_file_url(filename: str, folder: str) -> str:
    return f"https://{STORAGE_ACCOUNT_NAME}.file.core.windows.net/{FILE_SHARE_NAME}/{folder}/{filename}?{SAS_TOKEN}"

def retry_call(func_call, attempts: int = 3, delay: int = 1):
    last_exception = None
    for attempt in range(attempts):
        try:
            return func_call()
        except Exception as e:
            last_exception = e
            logging.error(f"Intento {attempt + 1} falló: {e}")
            time.sleep(delay)
    raise last_exception

def get_combined_analysis(kpr_desc: str, kr_desc: str, maturity_desc: str) -> tuple:
    """Envía un solo prompt con las descripciones de todas las gráficas y devuelve análisis individuales."""
    prompt = f"""
    Eres un analista experto en datos. Analiza estas gráficas y devuelve una interpretación analítica concisa para cada una:
    1. KPR: {kpr_desc}
    2. KR: {kr_desc}
    3. Madurez: {maturity_desc}
    Formato de respuesta: [KPR] texto [KR] texto [Madurez] texto
    Limita cada análisis a unas 20-30 palabras para optimizar espacio y tokens.
    """
    payload = {
        "messages": [
            {"role": "system", "content": "Proporciona análisis analíticos breves y precisos."},
            {"role": "user", "content": prompt}
        ],
        "max_tokens": 200,
        "temperature": 1.0,
        "top_p": 1.0
    }
    try:
        response = client.complete(payload)
        text = response.choices[0].message.content.strip()
    except Exception as e:
        logging.error(f"Error al invocar Azure OpenAI: {e}")
        return "Análisis no disponible.", "Análisis no disponible.", "Análisis no disponible."
    
    try:
        kpr = text.split("[KPR]")[1].split("[KR]")[0].strip()
        kr = text.split("[KR]")[1].split("[Madurez]")[0].strip()
        maturity = text.split("[Madurez]")[1].strip()
    except IndexError:
        logging.error("Formato de respuesta inválido, usando valores por defecto.")
        kpr = kr = maturity = "Análisis no disponible."
    return kpr, kr, maturity

def add_content_to_existing_slide(prs: Presentation, slide, images: list, analysis_text: str):
    """Agrega imágenes y análisis a una diapositiva existente."""
    margin = Inches(0.5)
    analysis_height = Inches(1)

    title_shape = None
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text.strip() and shape.top < Inches(2):
            title_shape = shape
            break

    title_bottom = title_shape.top + title_shape.height if title_shape else margin
    available_height = prs.slide_height - title_bottom - margin - analysis_height
    available_width = prs.slide_width - 2 * margin
    images_top = title_bottom
    images_left = margin

    num_images = len(images)
    if num_images == 0:
        logging.warning("No se proporcionaron imágenes para la diapositiva.")
        return

    if num_images == 1:
        rows, cols = 1, 1
    elif num_images <= 2:
        rows, cols = 1, num_images
    elif num_images <= 4:
        rows, cols = 2, 2
    else:
        cols = int(num_images ** 0.5)
        rows = cols if cols * cols >= num_images else cols + 1

    cell_width = available_width / cols
    cell_height = available_height / rows

    for i, img_stream in enumerate(images):
        col = i % cols
        row = i // cols
        cell_left = images_left + col * cell_width
        cell_top = images_top + row * cell_height

        img_stream.seek(0)
        img = Image.open(img_stream)
        orig_w, orig_h = img.size
        img_stream.seek(0)

        if (cell_width / cell_height) > (orig_w / orig_h):
            picture_height = cell_height
            picture_width = cell_height * (orig_w / orig_h)
        else:
            picture_width = cell_width
            picture_height = cell_width * (orig_h / orig_w)

        picture_left = cell_left + (cell_width - picture_width) / 2
        picture_top = cell_top + (cell_height - picture_height) / 2

        slide.shapes.add_picture(img_stream, picture_left, picture_top, width=picture_width, height=picture_height)

    analysisBox = slide.shapes.add_textbox(
        left=margin,
        top=prs.slide_height - margin - analysis_height,
        width=available_width,
        height=analysis_height
    )
    analysis_tf = analysisBox.text_frame
    analysis_tf.word_wrap = True  # Habilitar ajuste de texto
    analysis_tf.text = analysis_text
    for paragraph in analysis_tf.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(14)
            run.font.color.rgb = RGBColor(80, 80, 80)

def update_agenda_slide(prs: Presentation, slide, new_agenda_items: list):
    agenda_textbox = None
    for shape in slide.shapes:
        if shape.has_text_frame:
            text_frame = shape.text_frame
            if len(text_frame.paragraphs) > 1 and not text_frame.text.strip().lower().startswith("agenda"):
                agenda_textbox = shape
                break

    if not agenda_textbox:
        raise ValueError("No se encontró el cuadro de texto de la agenda en la diapositiva.")

    text_frame = agenda_textbox.text_frame
    original_paragraph = text_frame.paragraphs[0]
    original_font = original_paragraph.runs[0].font if original_paragraph.runs else None
    text_frame.clear()

    for item in new_agenda_items:
        p = text_frame.add_paragraph()
        p.text = item
        p.level = 0
        if original_font:
            run = p.runs[0]
            run.font.size = original_font.size
            run.font.name = original_font.name
            try:
                run.font.color.rgb = original_font.color.rgb
            except AttributeError:
                pass

def generate_kpr_graph(tmd_file_path: str, users_file_path: str, matricula_lider: str, q: str) -> tuple:
    df_users = pd.read_excel(users_file_path)
    df_tmd = pd.read_excel(tmd_file_path)
    users_list = df_users[df_users['Matricula Lider'] == matricula_lider]['Matricula'].tolist()
    df_filtered = df_tmd[df_tmd["Matricula LT"].isin(users_list)]
    plt.figure(figsize=(10, 6))
    plt.bar(df_filtered["Lider Técnico"], df_filtered[q], color='#000474')
    plt.title(f'{q} Values by Selected Lider Técnico (KPR)', fontsize=16)
    plt.xlabel('Lider Técnico', fontsize=12)
    plt.ylabel(q, fontsize=12)
    plt.xticks(rotation=45, ha='right')
    plt.tight_layout()
    img_bytes = BytesIO()
    plt.savefig(img_bytes, format='png')
    img_bytes.seek(0)
    plt.close()
    data_desc = f"Gráfica de {q} por Líder Técnico: {', '.join([f'{lt}: {val}' for lt, val in zip(df_filtered['Lider Técnico'], df_filtered[q])])}"
    return img_bytes, data_desc

def generate_kr_graph(pases_file_path: str, revisiones_file_path: str) -> tuple:
    quarters = ['Q1 01', 'Q1 02', 'Q1 03', 'Q2 01', 'Q2 02', 'Q2 03']
    df1 = pd.read_excel(pases_file_path)
    df2 = pd.read_excel(revisiones_file_path)
    df1_grouped = df1[quarters].sum(axis=0)
    df2_grouped = df2[quarters].sum(axis=0)
    plt.figure(figsize=(10, 7))
    x = range(len(quarters))
    plt.bar([pos - 0.2 for pos in x], df1_grouped, color='#000474', width=0.4, label='Pases Exitosos')
    plt.bar([pos + 0.2 for pos in x], df2_grouped, color='#f46a34', width=0.4, label='Reversiones')
    plt.title('Pases Exitosos vs Reversiones x Q (KR)', fontsize=16)
    plt.xlabel('Trimestre', fontsize=12)
    plt.ylabel('Valor Total', fontsize=12)
    plt.xticks(ticks=x, labels=quarters)
    plt.legend()
    plt.tight_layout()
    img_bytes = BytesIO()
    plt.savefig(img_bytes, format='png')
    img_bytes.seek(0)
    plt.close()
    data_desc = f"Pases Exitosos vs Reversiones por trimestre: Pases={dict(df1_grouped)}, Reversiones={dict(df2_grouped)}"
    return img_bytes, data_desc

def generate_maturity_graphs(maturity_file_path: str, users_file_path: str, matricula_lider: str, q: str) -> tuple:
    df_users = pd.read_excel(users_file_path)
    df_maturity = pd.read_excel(maturity_file_path)
    users_list = df_users[df_users['Matricula Lider'] == matricula_lider]['Matricula'].tolist()
    df_filtered = df_maturity[df_maturity["Matricula LT"].isin(users_list)]
    practices = df_filtered['PRACTICA'].unique()
    maturity_images = []
    maturity_desc_parts = []
    for practice in practices[:4]:
        df_practice = df_filtered[df_filtered['PRACTICA'] == practice]
        plt.figure(figsize=(10, 6))
        plt.bar(df_practice['Lider Técnico'], df_practice[q], color='#000474')
        plt.title(f'Niveles de Madurez para {practice}', fontsize=16)
        plt.xlabel('Lider Técnico', fontsize=12)
        plt.ylabel(f'Nivel de Madurez ({q})', fontsize=12)
        plt.xticks(rotation=45)
        plt.tight_layout()
        img_bytes = BytesIO()
        plt.savefig(img_bytes, format='png')
        img_bytes.seek(0)
        plt.close()
        maturity_images.append(img_bytes)
        maturity_desc_parts.append(f"{practice}: {', '.join([f'{lt}: {val}' for lt, val in zip(df_practice['Lider Técnico'], df_practice[q])])}")
    data_desc = "Niveles de madurez por práctica: " + "; ".join(maturity_desc_parts)
    return maturity_images, data_desc

app = func.FunctionApp(http_auth_level=func.AuthLevel.FUNCTION)

@app.function_name(name="generate_presentation")
@app.route(route="generate_presentation", methods=["POST"])
def generate_presentation(req: func.HttpRequest) -> func.HttpResponse:
    logging.info("Procesando solicitud de generate_presentation.")
    auth_header = req.headers.get("Authorization")
    if not auth_header or auth_header != AUTH_TOKEN_EXPECTED:
        return func.HttpResponse("No autorizado", status_code=401)
    
    try:
        req_body = req.get_json()
    except ValueError:
        return func.HttpResponse("JSON inválido en el cuerpo de la solicitud.", status_code=400)
    
    required_keys = ["q", "matricula_lider", "agenda"]
    missing_keys = [key for key in required_keys if key not in req_body]
    if missing_keys:
        return func.HttpResponse(f"Faltan claves requeridas: {missing_keys}", status_code=400)
    
    q = req_body["q"]
    matricula_lider = req_body["matricula_lider"]
    agenda_items = req_body["agenda"]
    
    with tempfile.TemporaryDirectory() as temp_dir:
        try:
            tmd_file_path = download_file_from_share(temp_dir, TMD_FILE, INPUTS_DIRECTORY)
            users_file_path = download_file_from_share(temp_dir, USERS_FILE, INPUTS_DIRECTORY)
            pases_file_path = download_file_from_share(temp_dir, PASES_FILE, INPUTS_DIRECTORY)
            revisiones_file_path = download_file_from_share(temp_dir, REVISIONES_FILE, INPUTS_DIRECTORY)
            maturity_file_path = download_file_from_share(temp_dir, MATURITY_LEVEL_FILE, INPUTS_DIRECTORY)
            template_file_path = download_file_from_share(temp_dir, "base_template.pptx", TEMPLATES_DIRECTORY)
            
            # Generar gráficas y descripciones
            kpr_img, kpr_desc = retry_call(lambda: generate_kpr_graph(tmd_file_path, users_file_path, matricula_lider, q))
            kr_img, kr_desc = retry_call(lambda: generate_kr_graph(pases_file_path, revisiones_file_path))
            maturity_imgs, maturity_desc = retry_call(lambda: generate_maturity_graphs(maturity_file_path, users_file_path, matricula_lider, q))
            
            # Obtener análisis combinado en una sola llamada
            kpr_analysis, kr_analysis, maturity_analysis = get_combined_analysis(kpr_desc, kr_desc, maturity_desc)
            
            prs = Presentation(template_file_path)
            if len(prs.slides) < 5:
                raise ValueError("La plantilla debe tener al menos 5 diapositivas.")
            
            update_agenda_slide(prs, prs.slides[1], agenda_items)
            add_content_to_existing_slide(prs, prs.slides[2], images=[kpr_img], analysis_text=kpr_analysis)
            add_content_to_existing_slide(prs, prs.slides[3], images=[kr_img], analysis_text=kr_analysis)
            add_content_to_existing_slide(prs, prs.slides[4], images=maturity_imgs, analysis_text=maturity_analysis)
            
            file_date = datetime.now().strftime("%m%d%y")
            output_filename = f"{file_date} Presentation.pptx"
            output_file_path = os.path.join(temp_dir, output_filename)
            prs.save(output_file_path)
            
            upload_file_to_share(output_file_path, output_filename, OUTPUTS_DIRECTORY)
            public_url = generate_file_url(output_filename, OUTPUTS_DIRECTORY)
            
            return func.HttpResponse(json.dumps({"public_url": public_url}), status_code=200, mimetype="application/json")
        
        except Exception as e:
            logging.error(f"Error en generate_presentation: {e}")
            return func.HttpResponse(f"Ocurrió un error: {str(e)}", status_code=500)